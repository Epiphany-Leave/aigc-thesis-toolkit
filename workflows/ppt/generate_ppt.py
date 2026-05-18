#!/usr/bin/env python3
"""Generate a defense PPT from a thesis document.

The PPT workflow is intentionally separate from the thesis writer.  It accepts
generated Markdown or externally supplied md/docx/pdf/txt files, uses a
PPT-specific AI planning flow, and writes independent outline/plan/preview
artifacts.
"""

from __future__ import annotations

import argparse
import json
import os
import re
import base64
import shlex
import shutil
import subprocess
import sys
import tempfile
import urllib.error
import urllib.request
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WORK = Path(__file__).resolve().parents[2]
CONFIG_FILE = WORK / "configs" / "default.yaml"
LOCAL_CONFIG_FILE = WORK / "configs" / "local.yaml"
OUTPUT_MD = WORK / "output" / "thesis.md"
OUTPUT_PPTX = WORK / "output" / "thesis_presentation.pptx"
SECTIONS_DIR = WORK / "thesis" / "sections"
PPT_DIR = WORK / "output" / "ppt"
PPT_PLAN = PPT_DIR / "plan.json"
PPT_OUTLINE = PPT_DIR / "outline.md"
PPT_PREVIEW = PPT_DIR / "preview.md"
PPT_VISUAL_DIR = PPT_DIR / "visuals"
PPT_SLIDE_IMAGE_DIR = PPT_DIR / "slide_images"
PPT_REFERENCE_STYLE = PPT_DIR / "reference_style.json"
PROMPT_FILE = WORK / "workflows" / "ppt" / "ppt_prompt.md"
MAX_SOURCE_CHARS = 90000
MAX_SLIDES = 16


THEMES = {
    "infographic": {
        "bg": RGBColor(248, 250, 252),
        "bg2": RGBColor(241, 245, 249),
        "panel": RGBColor(255, 255, 255),
        "panel2": RGBColor(235, 247, 255),
        "accent": RGBColor(25, 118, 210),
        "accent2": RGBColor(0, 150, 136),
        "accent3": RGBColor(251, 140, 0),
        "text": RGBColor(31, 41, 55),
        "muted": RGBColor(100, 116, 139),
    },
    "excalidraw": {
        "bg": RGBColor(255, 252, 242),
        "bg2": RGBColor(249, 244, 232),
        "panel": RGBColor(255, 255, 255),
        "panel2": RGBColor(255, 244, 220),
        "accent": RGBColor(48, 90, 176),
        "accent2": RGBColor(236, 116, 80),
        "accent3": RGBColor(29, 132, 108),
        "text": RGBColor(38, 38, 38),
        "muted": RGBColor(112, 112, 112),
    },
    "architecture": {
        "bg": RGBColor(245, 247, 250),
        "bg2": RGBColor(235, 241, 247),
        "panel": RGBColor(255, 255, 255),
        "panel2": RGBColor(232, 241, 248),
        "accent": RGBColor(47, 72, 88),
        "accent2": RGBColor(54, 162, 235),
        "accent3": RGBColor(0, 150, 136),
        "text": RGBColor(25, 33, 48),
        "muted": RGBColor(93, 105, 121),
    },
}


def clean_text(text: str) -> str:
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    text = re.sub(r"\[[^\]]+\]\([^)]+\)", "", text)
    text = re.sub(r"`[^`]*`", "", text)
    text = re.sub(r"[*_>#|]", "", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def compact_text(text: str, limit: int) -> str:
    text = clean_text(str(text))
    return text if len(text) <= limit else text[: max(0, limit - 3)].rstrip(",.; ") + "..."


def deep_merge(base, override):
    if not isinstance(base, dict) or not isinstance(override, dict):
        return override
    result = dict(base)
    for key, value in override.items():
        result[key] = deep_merge(result.get(key), value) if key in result else value
    return result


def load_config() -> dict:
    config = {}
    if CONFIG_FILE.exists():
        config = yaml.safe_load(CONFIG_FILE.read_text(encoding="utf-8")) or {}
    if LOCAL_CONFIG_FILE.exists():
        local = yaml.safe_load(LOCAL_CONFIG_FILE.read_text(encoding="utf-8")) or {}
        config = deep_merge(config, local)
    return config


def api_config(config: dict) -> tuple[str, str, str, int, float]:
    generation = config.get("engines", {}).get("generation", {})
    provider = generation.get("providers", {}).get("writer", {})
    base = (
        provider.get("api_base")
        or os.environ.get(provider.get("api_base_env", "OPENAI_BASE_URL"))
        or "https://api.openai.com/v1"
    ).rstrip("/")
    key = provider.get("api_key") or os.environ.get(provider.get("api_key_env", "OPENAI_API_KEY"), "")
    model = provider.get("model") or os.environ.get(provider.get("model_env", "OPENAI_MODEL"), "gpt-4o-mini")
    timeout = int(generation.get("batch", {}).get("request_timeout_seconds", 600))
    temperature = float(generation.get("ppt_temperature", 0.45))
    return base, key, model, timeout, temperature


def image_api_config(config: dict, image_model: str | None = None) -> tuple[str, str, str, str, int]:
    base, key, _model, timeout, _temperature = api_config(config)
    image_config = config.get("ppt", {}).get("image_slide", {})
    base = (
        image_config.get("api_base")
        or os.environ.get(image_config.get("api_base_env", "PPT_IMAGE_API_BASE"))
        or os.environ.get("OPENAI_IMAGE_BASE_URL")
        or base
    ).rstrip("/")
    key = (
        image_config.get("api_key")
        or os.environ.get(image_config.get("api_key_env", "PPT_IMAGE_API_KEY"))
        or os.environ.get("OPENAI_IMAGE_API_KEY")
        or key
    )
    model = (
        image_model
        or image_config.get("model")
        or os.environ.get("PPT_IMAGE_MODEL")
        or os.environ.get("OPENAI_IMAGE_MODEL")
        or "gpt-image-1"
    )
    size = image_config.get("size") or os.environ.get("PPT_IMAGE_SIZE") or "1536x1024"
    return base, key, model, size, int(image_config.get("timeout_seconds", timeout))


def visual_skill_config(config: dict) -> tuple[bool, list[str], int]:
    ppt_config = config.get("ppt", {}).get("visual_skill", {})
    command = ppt_config.get("command") or os.environ.get("PPT_VISUAL_SKILL_COMMAND", "")
    enabled = bool(ppt_config.get("enabled", bool(command)))
    timeout = int(ppt_config.get("timeout_seconds", 180))
    if isinstance(command, list):
        command_parts = [str(item) for item in command if str(item).strip()]
    else:
        command_parts = shlex.split(str(command), posix=os.name != "nt")
    if enabled and not command_parts:
        command_parts = [sys.executable, str(WORK / "workflows" / "ppt" / "visual_skill_runner.py")]
    return enabled and bool(command_parts), command_parts, timeout


def chat_completion(base: str, key: str, model: str, messages: list[dict], timeout: int, temperature: float) -> str:
    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": 7000,
        "response_format": {"type": "json_object"},
    }
    try:
        data = post_chat_completion(base, key, payload, timeout)
    except urllib.error.HTTPError as exc:
        if exc.code not in {400, 422}:
            raise
        payload.pop("response_format", None)
        data = post_chat_completion(base, key, payload, timeout)
    return data["choices"][0]["message"]["content"]


def post_chat_completion(base: str, key: str, payload: dict, timeout: int) -> dict:
    request = urllib.request.Request(
        f"{base}/chat/completions",
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {key}"},
        method="POST",
    )
    with urllib.request.urlopen(request, timeout=timeout) as response:
        return json.loads(response.read().decode("utf-8"))


def post_image_generation(base: str, key: str, payload: dict, timeout: int) -> dict:
    request = urllib.request.Request(
        f"{base}/images/generations",
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={"Content-Type": "application/json", "Authorization": f"Bearer {key}"},
        method="POST",
    )
    with urllib.request.urlopen(request, timeout=timeout) as response:
        return json.loads(response.read().decode("utf-8"))


def fetch_url_bytes(url: str, timeout: int) -> bytes:
    request = urllib.request.Request(url, headers={"User-Agent": "ppt-image-slide/1.0"})
    with urllib.request.urlopen(request, timeout=timeout) as response:
        return response.read()


def http_error_snippet(exc: urllib.error.HTTPError) -> str:
    try:
        body = exc.read().decode("utf-8", errors="replace")
    except Exception:
        body = ""
    body = clean_text(body)
    return body[:220]


def read_docx(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            if "word/document.xml" not in archive.namelist():
                return ""
            root = ET.fromstring(archive.read("word/document.xml"))
    except (OSError, zipfile.BadZipFile, ET.ParseError):
        return ""
    texts = [node.text or "" for node in root.iter() if node.tag.endswith("}t")]
    return clean_text("\n".join(texts))


def read_pdf(path: Path) -> str:
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", "-enc", "UTF-8", str(path), "-"],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=120,
        )
    except (OSError, subprocess.TimeoutExpired):
        return ""
    return clean_text(result.stdout) if result.returncode == 0 else ""


def convert_ppt_to_pptx(path: Path) -> Path | None:
    if path.suffix.lower() == ".pptx":
        return path
    if path.suffix.lower() != ".ppt" or shutil.which("libreoffice") is None:
        return None
    tmpdir = Path(tempfile.mkdtemp(prefix="ppt-template-"))
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pptx", "--outdir", str(tmpdir), str(path)],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=120,
        )
    except (OSError, subprocess.TimeoutExpired):
        return None
    converted = tmpdir / f"{path.stem}.pptx"
    return converted if result.returncode == 0 and converted.exists() else None


def shape_rgb(shape) -> str | None:
    try:
        fill = shape.fill
        if not fill or not fill.fore_color:
            return None
        color = fill.fore_color.rgb
        return rgb_to_hex(color) if color else None
    except Exception:
        return None


def percentile(values: list[float], ratio: float, fallback: float) -> float:
    if not values:
        return fallback
    values = sorted(values)
    index = min(len(values) - 1, max(0, round((len(values) - 1) * ratio)))
    return values[index]


def median(values: list[float], fallback: float) -> float:
    return percentile(values, 0.5, fallback)


def analyze_ppt_reference(path: Path) -> dict | None:
    converted = convert_ppt_to_pptx(path)
    if converted is None or not converted.exists():
        return None
    try:
        prs = Presentation(str(converted))
    except Exception as exc:
        print(f"PPT WARN: unable to analyze reference PPT: {path} {exc}")
        return None
    width = float(prs.slide_width or Inches(10))
    height = float(prs.slide_height or Inches(5.625))
    text_boxes: list[dict] = []
    media_boxes: list[dict] = []
    colors: list[str] = []
    layout_names: dict[str, int] = {}
    slide_count = len(prs.slides)
    for slide in prs.slides:
        layout_name = getattr(slide.slide_layout, "name", "") or "unknown"
        layout_names[layout_name] = layout_names.get(layout_name, 0) + 1
        for shape in slide.shapes:
            left = float(shape.left or 0) / width * 10
            top = float(shape.top or 0) / height * 5.625
            box_width = float(shape.width or 0) / width * 10
            box_height = float(shape.height or 0) / height * 5.625
            box = {"x": round(left, 3), "y": round(top, 3), "w": round(box_width, 3), "h": round(box_height, 3)}
            color = shape_rgb(shape)
            if color:
                colors.append(color)
            if getattr(shape, "has_text_frame", False):
                # Only geometry is stored; text content is intentionally ignored.
                text_boxes.append(box)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                media_boxes.append(box)
    text_lefts = [item["x"] for item in text_boxes if item["w"] > 0.25 and item["h"] > 0.12]
    text_tops = [item["y"] for item in text_boxes if item["w"] > 0.25 and item["h"] > 0.12]
    media_lefts = [item["x"] for item in media_boxes if item["w"] > 0.4 and item["h"] > 0.4]
    palette = []
    for color in colors:
        if color.lower() in {"#ffffff", "#000000"}:
            continue
        if color not in palette:
            palette.append(color)
        if len(palette) >= 6:
            break
    return {
        "source": str(path),
        "slides": slide_count,
        "slide_size": {"width": round(width, 2), "height": round(height, 2)},
        "layouts": layout_names,
        "text_box_count": len(text_boxes),
        "media_box_count": len(media_boxes),
        "text_anchor": {"x": round(median(text_lefts, 0.75), 3), "y": round(percentile(text_tops, 0.25, 1.35), 3)},
        "media_anchor": {"x": round(median(media_lefts, 6.35), 3)},
        "palette": palette,
    }


def analyze_ppt_references(paths: list[Path]) -> dict:
    profiles = [profile for path in paths if (profile := analyze_ppt_reference(path))]
    palettes = []
    for profile in profiles:
        for color in profile.get("palette", []):
            if color not in palettes:
                palettes.append(color)
    media_x = median([profile.get("media_anchor", {}).get("x", 6.35) for profile in profiles], 6.35)
    text_x = median([profile.get("text_anchor", {}).get("x", 0.75) for profile in profiles], 0.75)
    visual_side = "left" if media_x < text_x else "right"
    summary = {
        "sources": [profile["source"] for profile in profiles],
        "count": len(profiles),
        "palette": palettes[:6],
        "visual_side": visual_side,
        "text_anchor_x": round(text_x, 3),
        "media_anchor_x": round(media_x, 3),
        "profiles": profiles,
        "privacy": "Only layout geometry, master/layout names, shape categories and colors are stored. Text and image content are ignored.",
    }
    PPT_DIR.mkdir(parents=True, exist_ok=True)
    PPT_REFERENCE_STYLE.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    return summary


def apply_reference_style(theme: dict, reference_style: dict | None) -> dict:
    if not reference_style:
        return dict(theme)
    updated = dict(theme)
    palette = reference_style.get("palette") or []
    palette = [color for color in palette if usable_reference_accent(color)]
    if len(palette) >= 1:
        updated["accent"] = hex_to_rgb(palette[0], theme["accent"])
    if len(palette) >= 2:
        updated["accent2"] = hex_to_rgb(palette[1], theme["accent2"])
    if len(palette) >= 3:
        updated["accent3"] = hex_to_rgb(palette[2], theme["accent3"])
    return updated


def source_markdown() -> str:
    if OUTPUT_MD.exists():
        return OUTPUT_MD.read_text(encoding="utf-8-sig", errors="ignore")
    parts = []
    for path in sorted(SECTIONS_DIR.rglob("*.md")):
        content = path.read_text(encoding="utf-8-sig", errors="ignore").strip()
        if content:
            parts.append(content)
    if parts:
        return "\n\n".join(parts)
    raise SystemExit("ERROR: no thesis markdown found. Run python workflow.py build first, or pass --input.")


def read_source(path: Path | None) -> tuple[str, str]:
    if path is None:
        return source_markdown()[:MAX_SOURCE_CHARS], "output/thesis.md"
    if not path.exists():
        raise SystemExit(f"ERROR: input file not found: {path}")
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt", ".markdown"}:
        text = path.read_text(encoding="utf-8-sig", errors="ignore")
    elif suffix == ".docx":
        text = read_docx(path)
    elif suffix == ".pdf":
        text = read_pdf(path)
    else:
        text = path.read_text(encoding="utf-8-sig", errors="ignore")
    if not text.strip():
        raise SystemExit(f"ERROR: unable to extract text from {path}. Try converting it to md/docx/pdf text first.")
    return text[:MAX_SOURCE_CHARS], str(path)


def skip_heading(text: str) -> bool:
    compact = re.sub(r"\s+", "", text).lower()
    return bool(re.search(r"abstract|acknowledg|references|contents", compact))


def parse_headings(text: str) -> tuple[str, list[dict]]:
    title = ""
    chapters: list[dict] = []
    current: dict | None = None
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        heading = re.match(r"^(#{1,3})\s+(.+)$", line)
        chinese_heading = re.match(r"^Chapter\s+\d+\s+(.+)$", line, re.IGNORECASE)
        if heading or chinese_heading:
            heading_text = clean_text(heading.group(2) if heading else line)
            if not title and not skip_heading(heading_text):
                title = heading_text
            level = len(heading.group(1)) if heading else 1
            if level <= 2 and not skip_heading(heading_text) and heading_text != title:
                current = {"title": heading_text, "paragraphs": []}
                chapters.append(current)
            continue
        if current and not line.startswith(("$$", "|", "---")):
            cleaned = clean_text(line)
            if len(cleaned) >= 20:
                current["paragraphs"].append(cleaned)
    if not chapters:
        paragraphs = [clean_text(item) for item in re.split(r"\n{2,}", text) if len(clean_text(item)) > 30]
        chapters = [{"title": "Thesis core content", "paragraphs": paragraphs[:30]}]
    return title or chapters[0]["title"], chapters[:8]


def sentence_score(sentence: str) -> int:
    keywords = ["design", "implementation", "test", "result", "method", "system", "model", "control", "experiment", "analysis", "solution", "structure", "innovation"]
    return len(sentence) + sum(45 for key in keywords if key in sentence)


def bullets_from_paragraphs(paragraphs: list[str], limit: int = 4) -> list[str]:
    sentences: list[str] = []
    for paragraph in paragraphs[:20]:
        sentences.extend(item.strip() for item in re.split(r"[.;]", paragraph) if len(item.strip()) >= 12)
    ranked = sorted(sentences, key=sentence_score, reverse=True)
    bullets = []
    seen = set()
    for sentence in ranked:
        sentence = sentence[:80]
        key = sentence[:24]
        if key in seen:
            continue
        seen.add(key)
        bullets.append(sentence)
        if len(bullets) >= limit:
            break
    return bullets or ["Check this slide against the thesis body."]


def local_plan(text: str, style: str, source_name: str) -> dict:
    title, chapters = parse_headings(text)
    slides = [
        {
            "title": title,
            "kind": "cover",
            "layout": "cover",
            "bullets": ["Thesis defense presentation"],
            "visual_type": "hero",
            "visual": "Use the topic title as the main visual, formal and clear.",
            "notes": "Open with topic background and deck structure.",
        },
        {
            "title": "Agenda",
            "kind": "agenda",
            "layout": "agenda",
            "bullets": [chapter["title"] for chapter in chapters[:6]],
            "visual_type": "timeline",
            "visual": "Use a vertical agenda flow to show the presentation order.",
            "notes": "Briefly introduce the presentation order.",
        },
    ]
    for index, chapter in enumerate(chapters, start=1):
        bullets = bullets_from_paragraphs(chapter["paragraphs"])
        slides.append(
            {
                "title": f"{index}. {chapter['title']}",
                "kind": "content",
                "layout": "content_visual",
                "bullets": bullets,
                "visual_type": "process" if index % 2 else "architecture",
                "visual": "Draw a system structure, process, data chart, or image placeholder from this slide.",
                "diagram": bullets[:4],
                "notes": "Explain around the key points and avoid reading word by word.",
            }
        )
    slides.append(
        {
            "title": "Summary and Outlook",
            "kind": "summary",
            "layout": "summary",
            "bullets": ["Summarize system design and implementation", "Explain test conclusions and limitations", "Describe future optimization directions"],
            "visual_type": "summary",
            "visual": "Use a three-part conclusion diagram.",
            "diagram": ["Completed work", "Main conclusion", "Future work"],
            "notes": "Close with contributions and invite questions.",
        }
    )
    return {"title": title, "style": style, "source": source_name, "slides": slides[:14]}


def load_prompt() -> str:
    if PROMPT_FILE.exists():
        return PROMPT_FILE.read_text(encoding="utf-8")
    return "You are a thesis defense PPT director. Return strict JSON only."


def thesis_excerpt_for_slide(text: str, slide: dict) -> str:
    title = re.sub(r"^\d+[.]\s*", "", str(slide.get("title", ""))).strip()
    if not title:
        return text[:18000]
    lines = text.splitlines()
    hits = [idx for idx, line in enumerate(lines) if title[:12] and title[:12] in line]
    if not hits:
        return text[:18000]
    idx = hits[0]
    start = max(0, idx - 20)
    end = min(len(lines), idx + 180)
    excerpt = "\n".join(lines[start:end])
    return excerpt[:22000]


def parse_json_object(content: str) -> dict:
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", content, re.S)
        if not match:
            raise
        return json.loads(match.group(0))


def ai_outline(text: str, style: str, source_name: str, config: dict) -> dict | None:
    base, key, model, timeout, temperature = api_config(config)
    if not key:
        print("PPT INFO: no API key configured, using local planner.")
        return None
    prompt = load_prompt()
    user = (
        "Task: generate a global story line and slide blueprints first; do not write final dense content.\n"
        f"Visual preset: {style}\n"
        f"Input source: {source_name}\n"
        "layout must be one of cover, agenda, content_visual, visual_left, cards, statement, summary. Do not rotate layouts by page number.\n"
        "Return JSON: {title, style, source, narrative, slides:[{title, kind, layout, purpose, evidence_hint, visual_type, visual}]}\n"
        "Use 10-14 slides and cover title, agenda, background, solution/architecture, implementation, tests/results, summary.\n\n"
        f"Thesis content:\n{text[:MAX_SOURCE_CHARS]}"
    )
    try:
        print("PPT AI: generating deck outline...", flush=True)
        content = chat_completion(
            base,
            key,
            model,
            [{"role": "system", "content": prompt}, {"role": "user", "content": user}],
            timeout,
            temperature,
        )
        data = parse_json_object(content)
    except (json.JSONDecodeError, KeyError, urllib.error.URLError, TimeoutError, OSError) as exc:
        print(f"PPT WARN: AI outline failed, using local planner. {exc}")
        return None
    slides = data.get("slides") if isinstance(data, dict) else None
    if not isinstance(slides, list) or not slides:
        print("PPT WARN: AI outline returned no slides, using local planner.")
        return None
    data["title"] = str(data.get("title") or "Thesis defense presentation")
    data["style"] = style
    data["source"] = source_name
    data["slides"] = slides[:MAX_SLIDES]
    return data


def ai_refine_slide(outline: dict, slide: dict, index: int, total: int, text: str, config: dict) -> dict:
    base, key, model, timeout, temperature = api_config(config)
    prompt = load_prompt()
    excerpt = thesis_excerpt_for_slide(text, slide)
    user = (
        "Task: refine one PPT slide spec. Return only one JSON object.\n"
        f"Deck title: {outline.get('title', 'Thesis defense presentation')}\n"
        f"Slide: {index}/{total}\n"
        f"Visual preset: {outline.get('style', '')}\n"
        f"Narrative: {outline.get('narrative', '')}\n"
        f"Blueprint: {json.dumps(slide, ensure_ascii=False)}\n\n"
        "layout must be one of cover, agenda, content_visual, visual_left, cards, statement, summary.\n"
        "Return fields: title, kind, layout, bullets, visual_type, visual, diagram, callout, notes.\n"
        "Rules: bullets 3-5 items, each concise; diagram has 2-6 short drawable node labels; "
        "callout is one sentence; notes are speaker cues. Do not invent unsupported data.\n\n"
        f"Relevant thesis excerpt:\n{excerpt}"
    )
    print(f"PPT AI SLIDE: {index}/{total} {slide.get('title', '')}", flush=True)
    content = chat_completion(
        base,
        key,
        model,
        [{"role": "system", "content": prompt}, {"role": "user", "content": user}],
        timeout,
        temperature,
    )
    refined = parse_json_object(content)
    return {**slide, **refined}


def ai_plan(text: str, style: str, source_name: str, config: dict) -> dict | None:
    outline = ai_outline(text, style, source_name, config)
    if outline is None:
        return None
    base, key, _model, _timeout, _temperature = api_config(config)
    if not key or not base:
        return normalize_plan(outline)
    total = len(outline.get("slides", []))
    refined_slides = []
    for index, slide in enumerate(outline.get("slides", []), start=1):
        try:
            refined_slides.append(ai_refine_slide(outline, slide, index, total, text, config))
        except (json.JSONDecodeError, KeyError, urllib.error.URLError, TimeoutError, OSError) as exc:
            print(f"PPT WARN: slide {index} refinement failed, keeping outline slide. {exc}")
            refined_slides.append(slide)
    outline["slides"] = refined_slides
    return normalize_plan(outline)


def normalize_plan(plan: dict) -> dict:
    return {
        "title": clean_text(str(plan.get("title") or "Thesis defense presentation"))[:80],
        "style": str(plan.get("style") or "infographic"),
        "source": str(plan.get("source") or ""),
        "narrative": clean_text(str(plan.get("narrative") or ""))[:500],
        "slides": normalize_slides(plan.get("slides", [])),
    }


def normalize_slides(slides: list[dict]) -> list[dict]:
    normalized = []
    for slide in slides[:MAX_SLIDES]:
        if not isinstance(slide, dict):
            continue
        title = clean_text(str(slide.get("title", "")))[:60] or "Untitled slide"
        bullets = slide.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        bullets = [compact_text(item, 42) for item in bullets if clean_text(str(item))][:5]
        diagram = slide.get("diagram") or []
        if isinstance(diagram, str):
            diagram = [item.strip() for item in re.split(r"[,/;\n]", diagram) if item.strip()]
        diagram = [compact_text(item, 16) for item in diagram if clean_text(str(item))][:6]
        normalized.append(
            {
                "title": title,
                "kind": str(slide.get("kind") or "content"),
                "layout": str(slide.get("layout") or "content_visual"),
                "bullets": bullets or ["Extract this slide from the thesis body."],
                "visual_type": str(slide.get("visual_type") or slide.get("kind") or "process"),
                "visual": compact_text(slide.get("visual") or "Reserved visual area.", 140),
                "diagram": diagram,
                "callout": compact_text(slide.get("callout") or "", 58),
                "notes": clean_text(str(slide.get("notes") or ""))[:320],
            }
        )
    return normalized


def write_artifacts(plan: dict) -> None:
    PPT_DIR.mkdir(parents=True, exist_ok=True)
    PPT_PLAN.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")
    outline_lines = [f"# {plan.get('title', 'Thesis defense presentation')} PPT Outline", ""]
    preview_lines = [f"# {plan.get('title', 'Thesis defense presentation')} PPT Preview", ""]
    if plan.get("narrative"):
        outline_lines.extend(["## Narrative", plan["narrative"], ""])
    for index, slide in enumerate(plan.get("slides", []), start=1):
        outline_lines.append(f"{index}. {slide['title']} ({slide.get('kind', 'content')})")
        preview_lines.extend([f"## {index}. {slide['title']}", ""])
        preview_lines.extend(f"- {item}" for item in slide.get("bullets", []))
        if slide.get("callout"):
            preview_lines.extend(["", f"Callout: {slide['callout']}"])
        if slide.get("visual"):
            preview_lines.extend(["", f"Visual suggestion: {slide['visual']}"])
        if slide.get("notes"):
            preview_lines.extend(["", f"Speaker notes: {slide['notes']}"])
        preview_lines.append("")
    PPT_OUTLINE.write_text("\n".join(outline_lines).strip() + "\n", encoding="utf-8")
    PPT_PREVIEW.write_text("\n".join(preview_lines).strip() + "\n", encoding="utf-8")


def blank_presentation(reference_paths: list[Path] | None = None) -> Presentation:
    # Reference PPTs are design samples only. Starting from a blank package avoids
    # accidentally carrying over hidden media, old layouts, or embedded content.
    return Presentation()


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape(shape, fill: RGBColor, line: RGBColor | None = None, width: int = 1) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(width)


def rgb_to_hex(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def hex_to_rgb(value: str | None, fallback: RGBColor) -> RGBColor:
    if not value:
        return fallback
    value = value.strip().lstrip("#")
    if len(value) != 6:
        return fallback
    try:
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except ValueError:
        return fallback


def color_luminance(value: str | None) -> float:
    if not value:
        return 1.0
    value = value.strip().lstrip("#")
    if len(value) != 6:
        return 1.0
    try:
        r = int(value[0:2], 16) / 255
        g = int(value[2:4], 16) / 255
        b = int(value[4:6], 16) / 255
    except ValueError:
        return 1.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def usable_reference_accent(value: str | None) -> bool:
    if not value:
        return False
    clean = value.strip().lstrip("#")
    if len(clean) != 6:
        return False
    try:
        r, g, b = int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16)
    except ValueError:
        return False
    luminance = color_luminance(clean)
    channel_spread = max(r, g, b) - min(r, g, b)
    return luminance < 0.78 and channel_spread > 18


def slot(x: float, y: float, w: float, h: float) -> dict[str, float]:
    return {"x": x, "y": y, "w": w, "h": h}


def preferred_visual_left(reference_style: dict | None) -> bool:
    return bool(reference_style and reference_style.get("visual_side") == "left")


def slide_layout_name(item: dict, reference_style: dict | None) -> str:
    kind = str(item.get("kind") or "").lower()
    raw = str(item.get("layout") or "content_visual").lower()
    aliases = {
        "two_column": "content_visual",
        "content": "content_visual",
        "focus": "cards",
        "card": "cards",
    }
    layout = aliases.get(raw, raw)
    if kind == "agenda":
        return "agenda"
    if kind == "summary":
        return "summary"
    if layout not in {"agenda", "content_visual", "visual_left", "cards", "statement", "summary"}:
        layout = "content_visual"
    if layout == "content_visual" and preferred_visual_left(reference_style):
        return "visual_left"
    return layout


def layout_slots(layout: str) -> dict[str, dict[str, float]]:
    if layout == "visual_left":
        return {
            "body": slot(4.25, 1.45, 4.95, 3.1),
            "visual": slot(0.75, 1.45, 3.05, 3.25),
            "callout": slot(4.25, 4.72, 4.95, 0.42),
        }
    if layout == "summary":
        return {
            "body": slot(0.85, 1.45, 5.15, 2.95),
            "visual": slot(6.45, 1.65, 2.85, 2.65),
            "callout": slot(0.85, 4.72, 5.15, 0.42),
        }
    return {
        "body": slot(0.75, 1.45, 5.15, 3.15),
        "visual": slot(6.35, 1.45, 3.05, 3.25),
        "callout": slot(0.78, 4.72, 5.12, 0.42),
    }


def add_textbox(slide, left, top, width, height, text, size=24, color=None, bold=False, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def add_bullets(slide, bullets: list[str], theme, left=0.75, top=1.55, width=5.55, height=3.35, size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets[:5]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = theme["text"]
        paragraph.space_after = Pt(8)
    return box


def add_header(slide, text: str, theme, index: int, total: int) -> None:
    add_textbox(slide, Inches(0.55), Inches(0.32), Inches(7.85), Inches(0.55), text, 23, theme["text"], True)
    add_textbox(slide, Inches(8.55), Inches(0.38), Inches(0.9), Inches(0.3), f"{index:02d}/{total:02d}", 10, theme["muted"], False, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.04), Inches(2.35), Inches(0.05))
    set_shape(line, theme["accent"])


def add_footer(slide, theme, text: str = "Thesis defense") -> None:
    add_textbox(slide, Inches(0.55), Inches(5.28), Inches(4.2), Inches(0.18), text, 8, theme["muted"])


def add_decorative_band(slide, theme, variant: int) -> None:
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    set_shape(top, theme["bg2"])
    top.fill.transparency = 28
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.78), Inches(1.7), Inches(0.05))
    set_shape(line, theme["accent3" if variant % 2 else "accent"])


def add_callout(slide, text: str, theme, area: dict[str, float] | None = None) -> None:
    if not text:
        return
    text = compact_text(text, 58)
    area = area or slot(0.78, 4.72, 5.12, 0.42)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(area["x"]), Inches(area["y"]), Inches(area["w"]), Inches(area["h"]))
    set_shape(shape, theme["panel2"], theme["accent2"], 1)
    add_textbox(slide, Inches(area["x"] + 0.18), Inches(area["y"] + 0.09), Inches(area["w"] - 0.36), Inches(0.18), text, 10, theme["accent"], True, PP_ALIGN.CENTER)


def add_node(slide, x, y, w, h, text, theme, fill_key="panel", color_key="text", size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape(shape, theme[fill_key], theme["accent"], 1)
    add_textbox(slide, Inches(x + 0.08), Inches(y + 0.08), Inches(w - 0.16), Inches(h - 0.12), text, size, theme[color_key], True, PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2, theme):
    arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    arrow.line.color.rgb = theme["accent"]
    arrow.line.width = Pt(1.5)


def add_architecture_visual(slide, item: dict, theme, area: dict[str, float] | None = None) -> None:
    labels = item.get("diagram") or item.get("bullets", [])[:4]
    labels = (labels + ["Input", "Process", "Output"])[:4]
    area = area or slot(6.45, 1.45, 3.0, 3.35)
    left, top, width, height = area["x"], area["y"], area["w"], area["h"]
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape(panel, theme["panel"], theme["accent"], 1)
    add_textbox(slide, Inches(left + 0.18), Inches(top + 0.15), Inches(width - 0.36), Inches(0.35), "Structure", 12, theme["accent"], True, PP_ALIGN.CENTER)
    for idx, label in enumerate(labels[:4]):
        y = top + 0.65 + idx * min(0.62, max(0.48, (height - 1.1) / 4))
        add_node(slide, left + 0.35, y, max(1.8, width - 0.7), 0.38, label, theme, "panel2", "text", 10)
        if idx < min(len(labels), 4) - 1:
            add_arrow(slide, left + width / 2, y + 0.39, left + width / 2, y + 0.58, theme)


def add_process_visual(slide, item: dict, theme, area: dict[str, float] | None = None) -> None:
    labels = item.get("diagram") or item.get("bullets", [])[:4]
    labels = labels[:4] or ["Problem", "Solution", "Build", "Verify"]
    area = area or slot(6.35, 1.75, 3.05, 2.9)
    left, top, width = area["x"], area["y"], area["w"]
    for idx, label in enumerate(labels):
        y = top + idx * 0.72
        add_node(slide, left + 0.22, y, max(2.0, width - 0.6), 0.44, label, theme, "panel2", "text", 10)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left - 0.15), Inches(y + 0.05), Inches(0.34), Inches(0.34))
        set_shape(badge, theme["accent2" if idx % 2 else "accent"])
        add_textbox(slide, Inches(left - 0.09), Inches(y + 0.1), Inches(0.22), Inches(0.16), str(idx + 1), 8, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)


def add_compare_visual(slide, item: dict, theme, area: dict[str, float] | None = None) -> None:
    labels = item.get("diagram") or item.get("bullets", [])[:4]
    labels = labels[:4] or ["Current", "Improved", "Result"]
    area = area or slot(6.35, 1.55, 3.05, 2.9)
    left, top, width = area["x"], area["y"], area["w"]
    for idx, label in enumerate(labels[:4]):
        y = top + idx * 0.66
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.1), Inches(y + 0.28), Inches(max(1.2, width - 0.7 - idx * 0.28)), Inches(0.12))
        set_shape(bar, theme["accent2" if idx % 2 else "accent"])
        add_textbox(slide, Inches(left + 0.1), Inches(y), Inches(width - 0.25), Inches(0.25), label, 11, theme["text"], True)


def add_summary_visual(slide, item: dict, theme, area: dict[str, float] | None = None) -> None:
    labels = item.get("diagram") or item.get("bullets", [])[:3]
    labels = labels[:3] or ["Completed", "Conclusion", "Future work"]
    if area:
        left, top, width = area["x"], area["y"], area["w"]
        for idx, label in enumerate(labels):
            add_node(slide, left + 0.15, top + idx * 0.72, max(2.0, width - 0.3), 0.5, label, theme, "panel2", "text", 11)
        return
    for idx, label in enumerate(labels):
        add_node(slide, 1.0 + idx * 2.95, 2.85, 2.15, 0.82, label, theme, "panel2", "text", 13)


def add_agenda_cards(slide, bullets: list[str], theme) -> None:
    for idx, bullet in enumerate(bullets[:6]):
        row, col = divmod(idx, 2)
        x = 1.05 + col * 4.0
        y = 1.35 + row * 1.08
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.55), Inches(0.72))
        set_shape(card, theme["panel"], theme["panel2"], 1)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.36), Inches(0.36))
        set_shape(badge, theme["accent2" if idx % 2 else "accent"])
        add_textbox(slide, Inches(x + 0.25), Inches(y + 0.24), Inches(0.22), Inches(0.12), str(idx + 1), 8, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.68), Inches(y + 0.2), Inches(2.68), Inches(0.28), bullet, 13, theme["text"], True)


def add_focus_cards(slide, item: dict, theme) -> None:
    bullets = item.get("bullets", [])[:4]
    for idx, bullet in enumerate(bullets):
        x = 0.78 + (idx % 2) * 2.75
        y = 1.46 + (idx // 2) * 1.38
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.38), Inches(0.98))
        set_shape(card, theme["panel"], theme["panel2"], 1)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(0.98))
        set_shape(stripe, theme["accent2" if idx % 2 else "accent"])
        add_textbox(slide, Inches(x + 0.23), Inches(y + 0.18), Inches(1.96), Inches(0.52), compact_text(bullet, 28), 11, theme["text"], True, PP_ALIGN.CENTER)


def add_statement_layout(slide, item: dict, theme) -> None:
    callout = item.get("callout") or (item.get("bullets") or [""])[0]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(1.45), Inches(5.05), Inches(1.2))
    set_shape(shape, theme["panel2"], theme["accent"], 1)
    add_textbox(slide, Inches(1.08), Inches(1.75), Inches(4.45), Inches(0.5), compact_text(callout, 36), 17, theme["accent"], True, PP_ALIGN.CENTER)
    add_bullets(slide, item.get("bullets", [])[1:5], theme, left=0.95, top=3.0, width=4.75, height=1.55, size=13)


def visual_skill_payload(item: dict, index: int, total: int, style: str, theme: dict, reference_style: dict | None = None, area: dict[str, float] | None = None) -> dict:
    area = area or slot(6.35, 1.45, 3.05, 3.45)
    return {
        "contract": "ppt_visual_skill/v1",
        "index": index,
        "total": total,
        "style": style,
        "slide": item,
        "canvas": {"width_in": area["w"], "height_in": area["h"]},
        "theme": {key: rgb_to_hex(value) for key, value in theme.items()},
        "reference_style": reference_style or {},
        "output_dir": str(PPT_VISUAL_DIR),
    }


def run_visual_skill(item: dict, index: int, total: int, style: str, theme: dict, config: dict, reference_style: dict | None = None, area: dict[str, float] | None = None) -> dict | None:
    enabled, command, timeout = visual_skill_config(config)
    if not enabled:
        return None
    PPT_VISUAL_DIR.mkdir(parents=True, exist_ok=True)
    payload = visual_skill_payload(item, index, total, style, theme, reference_style, area)
    try:
        result = subprocess.run(
            command,
            input=json.dumps(payload, ensure_ascii=False),
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout,
            cwd=str(WORK),
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        print(f"PPT WARN: visual skill unavailable, using built-in visual. {exc}")
        return None
    if result.returncode != 0:
        stderr = clean_text(result.stderr)[-240:]
        print(f"PPT WARN: visual skill failed, using built-in visual. {stderr}")
        return None
    try:
        spec = parse_json_object(result.stdout)
    except json.JSONDecodeError as exc:
        print(f"PPT WARN: visual skill returned invalid JSON, using built-in visual. {exc}")
        return None
    return spec if isinstance(spec, dict) else None


def resolve_skill_path(path: str) -> Path:
    candidate = Path(path)
    return candidate if candidate.is_absolute() else WORK / candidate


def add_skill_image(slide, spec: dict, left: float, top: float, width: float, height: float) -> bool:
    image_path = spec.get("path") or spec.get("image") or spec.get("file")
    if not image_path:
        return False
    path = resolve_skill_path(str(image_path))
    if not path.exists():
        print(f"PPT WARN: visual skill image not found: {path}")
        return False
    try:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        return True
    except Exception as exc:
        print(f"PPT WARN: unable to insert visual skill image, using built-in visual. {exc}")
        return False


def add_skill_elements(slide, spec: dict, theme, left: float, top: float, width: float, height: float) -> bool:
    elements = spec.get("elements")
    if not isinstance(elements, list) or not elements:
        return False
    for element in elements[:24]:
        if not isinstance(element, dict):
            continue
        kind = str(element.get("type") or "box")
        x = left + float(element.get("x", 0))
        y = top + float(element.get("y", 0))
        w = max(0.05, float(element.get("w", element.get("width", 1))))
        h = max(0.05, float(element.get("h", element.get("height", 0.35))))
        text = clean_text(str(element.get("text", "")))
        fill = hex_to_rgb(str(element.get("fill", "")), theme["panel2"])
        line = hex_to_rgb(str(element.get("line", "")), theme["accent"])
        color = hex_to_rgb(str(element.get("color", "")), theme["text"])
        if kind == "image":
            add_skill_image(slide, element, x, y, w, h)
        elif kind == "text":
            add_textbox(slide, Inches(x), Inches(y), Inches(w), Inches(h), text, int(element.get("size", 11)), color, bool(element.get("bold", False)), PP_ALIGN.CENTER if element.get("align") == "center" else None)
        elif kind == "arrow":
            add_arrow(slide, x, y, left + float(element.get("x2", element.get("to_x", 0))), top + float(element.get("y2", element.get("to_y", 0))), theme)
        elif kind == "bar":
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            set_shape(shape, fill, line, int(element.get("line_width", 1)))
        else:
            shape_type = MSO_SHAPE.OVAL if kind in {"circle", "oval"} else MSO_SHAPE.ROUNDED_RECTANGLE
            shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
            set_shape(shape, fill, line, int(element.get("line_width", 1)))
            if text:
                add_textbox(slide, Inches(x + 0.04), Inches(y + 0.06), Inches(max(0.05, w - 0.08)), Inches(max(0.05, h - 0.1)), text, int(element.get("size", 10)), color, bool(element.get("bold", True)), PP_ALIGN.CENTER)
    return True


def add_skill_visual(slide, item: dict, theme, style: str, config: dict, index: int, total: int, reference_style: dict | None = None, area: dict[str, float] | None = None) -> bool:
    area = area or slot(6.35, 1.45, 3.05, 3.45)
    spec = run_visual_skill(item, index, total, style, theme, config, reference_style, area)
    if not spec:
        return False
    left, top, width, height = area["x"], area["y"], area["w"], area["h"]
    if str(spec.get("type", "")).lower() == "image" and add_skill_image(slide, spec, left, top, width, height):
        return True
    if add_skill_elements(slide, spec, theme, left, top, width, height):
        return True
    if add_skill_image(slide, spec, left, top, width, height):
        return True
    print("PPT WARN: visual skill output had no usable image/elements, using built-in visual.")
    return False


def add_visual(slide, item: dict, theme, style: str, config: dict, index: int, total: int, reference_style: dict | None = None, area: dict[str, float] | None = None) -> None:
    area = area or slot(6.35, 1.45, 3.05, 3.25)
    if add_skill_visual(slide, item, theme, style, config, index, total, reference_style, area):
        return
    visual_type = (item.get("visual_type") or "").lower()
    if "arch" in visual_type or item.get("kind") == "architecture" or style == "architecture":
        add_architecture_visual(slide, item, theme, area)
    elif "compare" in visual_type or "result" in item.get("kind", ""):
        add_compare_visual(slide, item, theme, area)
    elif "summary" in visual_type or item.get("kind") == "summary":
        add_summary_visual(slide, item, theme, area)
    else:
        add_process_visual(slide, item, theme, area)


def reference_style_prompt(reference_style: dict | None) -> str:
    if not reference_style:
        return "No external reference deck was provided."
    palette = ", ".join(reference_style.get("palette") or []) or "use a safe high-contrast academic palette"
    side = reference_style.get("visual_side") or "right"
    return (
        "Reference deck style summary: borrow only visual style, layout rhythm, color tendency, "
        f"and image/text placement. Never reuse old text or pictures. Palette candidates: {palette}. "
        f"Typical visual side: {side}."
    )


def slide_image_prompt(item: dict, index: int, total: int, style: str, reference_style: dict | None) -> str:
    bullets = "\n".join(f"- {bullet}" for bullet in item.get("bullets", [])[:5])
    diagram = ", ".join(item.get("diagram", [])[:6])
    return (
        "Create one complete 16:9 thesis defense PowerPoint slide as a single polished slide image. Keep all important content inside a 16:9 safe area.\n"
        "The image must include all visible content for this slide: title, key points, diagram, background, shapes, and layout.\n"
        "Use large readable text, strong contrast, clean spacing, and avoid overlap. Keep Chinese text minimal and legible.\n"
        "Do not create a poster, landing page, phone UI, mockup frame, watermark, logo, or browser chrome.\n"
        "Do not invent experimental data, formulas, citations, or chart values not provided here.\n"
        f"Visual preset: {style}.\n"
        f"{reference_style_prompt(reference_style)}\n"
        f"Slide {index}/{total}\n"
        f"Title: {item.get('title', '')}\n"
        f"Layout intent: {item.get('layout', 'content_visual')}\n"
        f"Visual type: {item.get('visual_type', 'process')}\n"
        f"Callout: {item.get('callout', '')}\n"
        f"Bullets:\n{bullets}\n"
        f"Diagram node labels: {diagram}\n"
        f"Visual instruction: {item.get('visual', '')}\n"
        "Output should look like a finished academic defense slide, with a restrained professional style."
    )


def decode_image_response(data: dict, timeout: int) -> bytes | None:
    items = data.get("data") if isinstance(data, dict) else None
    if not isinstance(items, list) or not items:
        return None
    first = items[0]
    if not isinstance(first, dict):
        return None
    b64 = first.get("b64_json") or first.get("base64")
    if b64:
        return base64.b64decode(b64)
    url = first.get("url")
    if url:
        return fetch_url_bytes(str(url), timeout)
    return None


def render_slide_image(item: dict, index: int, total: int, style: str, config: dict, reference_style: dict | None, image_model: str | None = None) -> Path | None:
    base, key, model, size, timeout = image_api_config(config, image_model)
    if not key:
        print("PPT WARN: no API key configured for image_slide mode.")
        return None
    PPT_SLIDE_IMAGE_DIR.mkdir(parents=True, exist_ok=True)
    prompt = slide_image_prompt(item, index, total, style, reference_style)
    payload = {
        "model": model,
        "prompt": prompt,
        "size": size,
        "n": 1,
        "response_format": "b64_json",
    }
    try:
        data = post_image_generation(base, key, payload, timeout)
    except urllib.error.HTTPError as exc:
        if exc.code in {400, 422}:
            payload.pop("response_format", None)
            try:
                data = post_image_generation(base, key, payload, timeout)
            except urllib.error.HTTPError as retry_exc:
                print(f"PPT WARN: image generation API failed for slide {index}: HTTP {retry_exc.code} {http_error_snippet(retry_exc)}")
                return None
        else:
            print(f"PPT WARN: image generation endpoint unavailable for slide {index}: HTTP {exc.code} {http_error_snippet(exc)}")
            print("PPT WARN: configure ppt.image_slide.api_base/api_key/model or PPT_IMAGE_API_BASE/PPT_IMAGE_API_KEY/PPT_IMAGE_MODEL for image_slide mode.")
            return None
    except (urllib.error.URLError, TimeoutError, OSError) as exc:
        print(f"PPT WARN: image generation request failed for slide {index}: {exc}")
        return None
    image_bytes = decode_image_response(data, timeout)
    if not image_bytes:
        print(f"PPT WARN: image generation returned no image for slide {index}.")
        return None
    path = PPT_SLIDE_IMAGE_DIR / f"slide_{index:02d}.png"
    path.write_bytes(image_bytes)
    return path


def add_notes(slide, text: str) -> None:
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def build_image_slide_presentation(plan: dict, style: str, reference_paths: list[Path] | None = None, config: dict | None = None, reference_style: dict | None = None, image_model: str | None = None) -> Presentation | None:
    config = config or {}
    slides = normalize_slides(plan.get("slides", []))
    prs = blank_presentation(reference_paths)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    total = max(1, len(slides))
    print(f"PPT TOTAL: {total}")
    for index, item in enumerate(slides, start=1):
        print(f"PPT IMAGE: {index}/{total} {item['title']}", flush=True)
        image_path = render_slide_image(item, index, total, style, config, reference_style, image_model)
        if image_path is None:
            return None
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        add_notes(slide, item.get("notes", ""))
    return prs


def build_presentation(plan: dict, style: str, reference_paths: list[Path] | None = None, config: dict | None = None, reference_style: dict | None = None) -> Presentation:
    config = config or {}
    theme = apply_reference_style(THEMES.get(style, THEMES["infographic"]), reference_style)
    slides = normalize_slides(plan.get("slides", []))
    prs = blank_presentation(reference_paths)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    total = max(1, len(slides))
    print(f"PPT TOTAL: {total}")
    for index, item in enumerate(slides, start=1):
        print(f"PPT PROGRESS: {index}/{total} {item['title']}", flush=True)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, theme["bg"])
        kind = item.get("kind", "content")
        layout = slide_layout_name(item, reference_style)
        areas = layout_slots(layout)
        if kind == "cover" or index == 1:
            hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
            set_shape(hero, theme["bg2"])
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.55), Inches(5.625))
            set_shape(accent, theme["accent"])
            ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(4.82), Inches(9.45), Inches(0.34))
            set_shape(ribbon, theme["accent2"])
            add_textbox(slide, Inches(0.95), Inches(1.18), Inches(7.95), Inches(1.25), item["title"], 30, theme["text"], True, PP_ALIGN.LEFT)
            subtitle = item["bullets"][0] if item.get("bullets") else "Thesis defense presentation"
            add_textbox(slide, Inches(1.0), Inches(2.62), Inches(5.8), Inches(0.42), subtitle, 17, theme["muted"], False, PP_ALIGN.LEFT)
            add_summary_visual(slide, {"diagram": ["Background", "Design", "Testing"]}, theme)
            continue
        add_decorative_band(slide, theme, index)
        add_header(slide, item["title"], theme, index, total)
        if kind == "agenda" or layout == "agenda":
            add_agenda_cards(slide, item.get("bullets", []), theme)
        elif kind == "summary" or layout == "summary":
            body = areas["body"]
            add_bullets(slide, item.get("bullets", []), theme, left=body["x"], top=body["y"], width=body["w"], height=body["h"], size=15)
            if not add_skill_visual(slide, item, theme, style, config, index, total, reference_style, areas["visual"]):
                add_summary_visual(slide, item, theme, areas["visual"])
        elif layout == "statement":
            add_statement_layout(slide, item, theme)
            add_visual(slide, item, theme, style, config, index, total, reference_style, areas["visual"])
        elif layout == "cards":
            add_focus_cards(slide, item, theme)
            add_visual(slide, item, theme, style, config, index, total, reference_style, areas["visual"])
        else:
            body = areas["body"]
            add_bullets(slide, item.get("bullets", []), theme, left=body["x"], top=body["y"], width=body["w"], height=body["h"], size=16)
            add_visual(slide, item, theme, style, config, index, total, reference_style, areas["visual"])
        if layout != "agenda":
            add_callout(slide, item.get("callout", ""), theme, areas.get("callout"))
        add_footer(slide, theme)
    return prs


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", help="input thesis file: md/docx/pdf/txt. Default: output/thesis.md")
    parser.add_argument("--output", default=str(OUTPUT_PPTX), help="output pptx path")
    parser.add_argument("--style", default="infographic", choices=sorted(THEMES), help="visual style preset")
    parser.add_argument("--template", action="append", help="optional reference PPT design sample: pptx or ppt. Can be used multiple times")
    parser.add_argument("--render-mode", default="editable", choices=["editable", "image_slide"], help="PPT rendering mode: editable shapes or full-slide AI images")
    parser.add_argument("--image-model", help="image generation model for --render-mode image_slide")
    parser.add_argument("--no-ai", action="store_true", help="skip AI planning and use local extraction")
    args = parser.parse_args()

    source_path = Path(args.input) if args.input else None
    text, source_name = read_source(source_path)
    config = load_config()
    plan = None if args.no_ai else ai_plan(text, args.style, source_name, config)
    if plan is None:
        plan = normalize_plan(local_plan(text, args.style, source_name))
    else:
        plan = normalize_plan(plan)
    write_artifacts(plan)
    reference_paths = [Path(item) for item in (args.template or [])]
    if reference_paths:
        reference_style = analyze_ppt_references(reference_paths)
    else:
        reference_style = None
        if PPT_REFERENCE_STYLE.exists():
            PPT_REFERENCE_STYLE.unlink()
    if args.render_mode == "image_slide":
        prs = build_image_slide_presentation(plan, args.style, reference_paths=reference_paths, config=config, reference_style=reference_style, image_model=args.image_model)
        if prs is None:
            raise SystemExit(
                "ERROR: image_slide mode failed. No editable fallback was generated. "
                "Configure PPT 图像 API in WebUI, or switch PPT 渲染模式 to 可编辑 PPT 元素."
            )
    else:
        prs = build_presentation(plan, args.style, reference_paths=reference_paths, config=config, reference_style=reference_style)
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"OK: pptx -> {output}")
    print(f"OK: outline -> {PPT_OUTLINE}")
    print(f"OK: preview -> {PPT_PREVIEW}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
